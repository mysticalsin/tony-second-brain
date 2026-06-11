#!/usr/bin/env python3
"""Convert binary documents (PDF/Word/PowerPoint/Excel) across the vault and its
_External linked folders into markdown notes inside the vault.

- Sources: vault real files + each _External/<name> symlink target (followed).
- Output:  "<vault>/Document Library/<source-label>/<relpath>.md" (mirrors tree).
- Privacy: skips any path containing "HR Documents" or "LinkedIn" (matches the
  capture/ingest guards in CLAUDE.md).
- Incremental: skips a file whose markdown is newer than the source (unless --force).
- Robust:  per-file try/except; failures logged, run continues (fail-visibly).
- Pipeline-friendly: --max-seconds stops gracefully so hourly runs stay cheap;
  the first full run processes everything, later runs only touch new/changed docs.

Wired as a step in 99_Meta/brain-refresh.sh.
"""
from __future__ import annotations

import argparse
import csv
import datetime as dt
import os
import re
import subprocess
import sys
import time
from pathlib import Path

VAULT = Path(os.environ.get("VAULT_ROOT") or (_ for _ in ()).throw(SystemExit("Set VAULT_ROOT to your vault path")))
OUT_DIRNAME = "Document Library"
MARKITDOWN = os.path.expanduser("~/.local/bin/markitdown")

EXTS = {".pdf", ".docx", ".doc", ".pptx", ".ppt", ".xlsx", ".xls"}
PRIVACY = ("hr documents", "linkedin", "meetings/confidential")  # lowercase substrings to skip
# directories never descended into when walking the vault itself
VAULT_PRUNE = {
    ".obsidian", "_External", "_brain_index", "_brain_api", "out",
    "graphify-out", "build", "node_modules", ".git", OUT_DIRNAME,
}


def is_private(path: str) -> bool:
    p = path.lower()
    return any(frag in p for frag in PRIVACY)


def slug_title(name: str) -> str:
    return re.sub(r"\s+", " ", name).strip()


def iter_sources():
    """Yield (label, root_path) pairs: the vault itself, then each _External link."""
    yield ("vault", VAULT)
    ext = VAULT / "_External"
    if ext.is_dir():
        for entry in sorted(ext.iterdir()):
            try:
                if entry.is_symlink() and entry.resolve().is_dir():
                    yield (entry.name, entry.resolve())
            except OSError:
                continue


def iter_docs(label: str, root: Path):
    """Walk a source root, yielding (source_file, rel_path) for convertible docs."""
    follow = label != "vault"  # external roots are symlinks → follow; vault is real
    for dirpath, dirnames, filenames in os.walk(root, followlinks=follow):
        if label == "vault":
            dirnames[:] = [d for d in dirnames if d not in VAULT_PRUNE]
        # never descend into private trees regardless of source
        dirnames[:] = [d for d in dirnames if not is_private(d)]
        for fn in filenames:
            if Path(fn).suffix.lower() not in EXTS:
                continue
            src = Path(dirpath) / fn
            if is_private(str(src)):
                continue
            try:
                rel = src.relative_to(root)
            except ValueError:
                rel = Path(fn)
            yield src, rel


def target_for(label: str, rel: Path) -> Path:
    return VAULT / OUT_DIRNAME / label / rel.with_suffix(rel.suffix + ".md")


def needs_convert(src: Path, dst: Path, force: bool) -> bool:
    if force or not dst.exists():
        return True
    try:
        return src.stat().st_mtime > dst.stat().st_mtime
    except OSError:
        return True


def _fallback_extract(src: Path) -> str:
    """Recover text when markitdown fails.

    - .pptx: python-pptx (markitdown throws 'no embedded image' on image-less
      slides — a known bug; the text is still there).
    - .doc/.rtf/.docx: macOS `textutil` (native, handles legacy binary .doc).
    Returns extracted text, or '' if the fallback can't help either.
    """
    suf = src.suffix.lower()
    if suf == ".pptx":
        try:
            from pptx import Presentation
            prs = Presentation(str(src))
            out = []
            for i, slide in enumerate(prs.slides, 1):
                out.append(f"\n## Slide {i}\n")
                for shape in slide.shapes:
                    if shape.has_text_frame and shape.text_frame.text.strip():
                        out.append(shape.text_frame.text.strip())
                    if shape.has_table:
                        for row in shape.table.rows:
                            out.append(" | ".join(c.text for c in row.cells))
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    note = slide.notes_slide.notes_text_frame.text.strip()
                    if note:
                        out.append(f"\n_Notes:_ {note}")
            return "\n\n".join(out).strip()
        except Exception:
            return ""
    if suf in (".doc", ".rtf", ".docx"):
        try:
            r = subprocess.run(["textutil", "-convert", "txt", "-stdout", str(src)],
                               capture_output=True, text=True, timeout=120)
            return r.stdout.strip() if r.returncode == 0 else ""
        except (OSError, subprocess.TimeoutExpired):
            return ""
    return ""


def convert_one(src: Path, dst: Path, label: str) -> tuple[bool, str]:
    body = ""
    used_fallback = False
    try:
        proc = subprocess.run(
            [MARKITDOWN, str(src)],
            capture_output=True, text=True, timeout=180,
        )
        if proc.returncode == 0:
            body = proc.stdout.strip()
    except subprocess.TimeoutExpired:
        pass  # try fallback below
    except OSError as e:
        return False, f"spawn: {e}"

    if not body:
        body = _fallback_extract(src)
        used_fallback = bool(body)
        if not body:
            return False, "markitdown + fallback both failed (corrupt / unsupported)"

    now = dt.datetime.now().astimezone().isoformat(timespec="seconds")
    fm = (
        "---\n"
        "type: converted-doc\n"
        f"title: \"{slug_title(src.stem)[:120].replace(chr(34), chr(39))}\"\n"
        f"source_path: \"{src}\"\n"
        f"source_format: {src.suffix.lower().lstrip('.')}\n"
        f"source_collection: {label}\n"
        f"converted: {now}\n"
        f"converter: {'fallback' if used_fallback else 'markitdown'}\n"
        "tags: [converted-doc]\n"
        "---\n\n"
        f"> Auto-converted from `{src.name}` ({label}). Regenerated by "
        "`build/tools/convert_docs_to_md.py`. Original is the source of truth.\n\n"
    )
    dst.parent.mkdir(parents=True, exist_ok=True)
    tmp = dst.with_suffix(dst.suffix + ".tmp")
    tmp.write_text(fm + body, encoding="utf-8")
    os.replace(tmp, dst)
    return True, "ok"


def write_manifest():
    """Central traceability: one row per source doc → 'Document Library/_manifest.csv'.
    Records source, target, the converter actually used (markitdown / fallback) and
    status. Pure filesystem read — no LLM, no tokens. Refreshed every run."""
    rows = []
    for label, root in iter_sources():
        for src, rel in iter_docs(label, root):
            dst = target_for(label, rel)
            converter = converted_at = ""
            if dst.exists():
                status = "converted"
                try:
                    head = dst.read_text(encoding="utf-8", errors="replace")[:800]
                    m = re.search(r'^converter:\s*(\w+)', head, re.M); converter = m.group(1) if m else "?"
                    m2 = re.search(r'^converted:\s*(\S+)', head, re.M); converted_at = m2.group(1) if m2 else ""
                except OSError:
                    status = "unreadable"
            else:
                status = "not-converted"
            try: sbytes = src.stat().st_size
            except OSError: sbytes = 0
            rows.append([str(src), label, src.suffix.lower().lstrip('.'), sbytes, str(dst), status, converter, converted_at])
    out = VAULT / "Document Library" / "_manifest.csv"
    out.parent.mkdir(parents=True, exist_ok=True)
    with out.open("w", newline="", encoding="utf-8") as f:
        w = csv.writer(f)
        w.writerow(["source_path", "collection", "format", "source_bytes", "target", "status", "converter", "converted_at"])
        w.writerows(sorted(rows))
    conv = sum(1 for r in rows if r[5] == "converted")
    return out, len(rows), conv


def main() -> int:
    ap = argparse.ArgumentParser()
    ap.add_argument("--force", action="store_true", help="reconvert even if up to date")
    ap.add_argument("--limit", type=int, default=0, help="max files to convert (0=all)")
    ap.add_argument("--max-seconds", type=int, default=0, help="stop gracefully after N seconds (0=no limit)")
    ap.add_argument("--quiet", action="store_true")
    args = ap.parse_args()

    if not Path(MARKITDOWN).exists():
        print(f"FATAL: markitdown not found at {MARKITDOWN}", file=sys.stderr)
        return 2

    start = time.time()
    converted = skipped = failed = 0
    failures: list[str] = []

    for label, root in iter_sources():
        for src, rel in iter_docs(label, root):
            if args.max_seconds and (time.time() - start) > args.max_seconds:
                if not args.quiet:
                    print(f"[budget] stopping after {args.max_seconds}s")
                _summary(converted, skipped, failed, failures, start, partial=True)
                return 0
            if args.limit and converted >= args.limit:
                _summary(converted, skipped, failed, failures, start, partial=True)
                return 0

            dst = target_for(label, rel)
            if not needs_convert(src, dst, args.force):
                skipped += 1
                continue
            ok, msg = convert_one(src, dst, label)
            if ok:
                converted += 1
                if not args.quiet and converted % 25 == 0:
                    print(f"  …{converted} converted ({int(time.time()-start)}s)")
            else:
                failed += 1
                failures.append(f"{src} :: {msg}")

    mpath, mtotal, mconv = write_manifest()
    if not args.quiet:
        print(f"[manifest] {mconv}/{mtotal} docs traced → {mpath}")
    _summary(converted, skipped, failed, failures, start, partial=False)
    return 0


def _summary(conv, skip, fail, failures, start, partial):
    tag = "partial" if partial else "complete"
    print(f"[convert_docs] {tag}: converted={conv} skipped={skip} failed={fail} "
          f"in {int(time.time()-start)}s → {VAULT / OUT_DIRNAME}")
    if failures:
        log = VAULT / "99_Meta" / "convert-docs-errors.log"
        try:
            log.write_text("\n".join(failures) + "\n", encoding="utf-8")
            print(f"[convert_docs] {len(failures)} failures logged → {log}")
        except OSError:
            for f in failures[:10]:
                print("  FAIL", f)


if __name__ == "__main__":
    sys.exit(main())
